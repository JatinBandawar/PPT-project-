import cv2
import numpy as np
import speech_recognition as sr
from pptx import Presentation

# Initialize PowerPoint presentation
presentation_path = "D:\\Snehal\\PBL\\PPT NAVIGATOR\\New Microsoft PowerPoint Presentation.pptx"
prs = Presentation(presentation_path)
slide_count = len(prs.slides)
current_slide_index = 0

# Function to move to the next slide
def next_slide():
    global current_slide_index
    if current_slide_index < slide_count - 1:
        current_slide_index += 1
        prs.slides[current_slide_index].visible = True

# Function to move to the previous slide
def previous_slide():
    global current_slide_index
    if current_slide_index > 0:
        prs.slides[current_slide_index].visible = False
        current_slide_index -= 1

# Function to perform action based on voice command
def perform_action(command):
    if "next" in command:
        next_slide()
    elif "previous" in command:
        previous_slide()
    # Add more commands as needed

# Function to listen for voice commands
def listen():
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        print("Listening for commands...")
        recognizer.adjust_for_ambient_noise(source)
        audio = recognizer.listen(source)
    try:
        command = recognizer.recognize_google(audio).lower()
        print("You said:", command)
        perform_action(command)
    except sr.UnknownValueError:
        print("Could not understand audio")
    except sr.RequestError as e:
        print("Could not request results; {0}".format(e))

# Function to detect hand gestures using OpenCV
def detect_gesture():
    cap = cv2.VideoCapture(0)
    while True:
        ret, frame = cap.read()
        # Add your hand gesture detection logic here
        # This could involve contour detection, hand tracking, etc.
        # For simplicity, I'm leaving this part empty
        cv2.imshow('Frame', frame)
        if cv2.waitKey(1) & 0xFF == ord('q'):
            break
    cap.release()
    cv2.destroyAllWindows()

# Main loop to continuously listen for commands and detect gestures
while True:
    listen()
    # Optionally, you can uncomment the line below to detect gestures using OpenCV
    # detect_gesture()